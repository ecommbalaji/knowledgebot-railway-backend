#!/usr/bin/env python3
"""
Generate binary test files (PDF, DOCX, PPTX, XLSX, PNG, JPEG) with current affairs Q&A content.
Requires: pip install reportlab python-docx python-pptx openpyxl pillow
"""

import os
import sys

# Q&A Content
QA_PAIRS = [
    ("What is the capital of Australia?", "Canberra is the capital of Australia."),
    ("Who is the UN Secretary-General?", "António Guterres from Portugal since 2017."),
    ("What is the Paris Agreement?", "International climate treaty adopted in 2015."),
    ("Which country has the largest population?", "India surpassed China in 2023."),
    ("What is cryptocurrency?", "Digital currency using cryptography."),
    ("What are the BRICS nations?", "Brazil, Russia, India, China, South Africa + more."),
    ("What is AI?", "Simulation of human intelligence by machines."),
    ("What is the World Economic Forum?", "International org hosting Davos meetings."),
    ("What is COP?", "Annual UN climate change conferences."),
    ("What is the ISS?", "Modular space station in low Earth orbit."),
]

def create_pdf():
    """Create a PDF file using reportlab."""
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
        from reportlab.lib.units import inch
        
        filename = "current_affairs_qa.pdf"
        c = canvas.Canvas(filename, pagesize=letter)
        width, height = letter
        
        # Title
        c.setFont("Helvetica-Bold", 18)
        c.drawString(1*inch, height - 1*inch, "Current Affairs Q&A 2026")
        
        # Q&A pairs
        c.setFont("Helvetica", 11)
        y = height - 1.5*inch
        
        for i, (q, a) in enumerate(QA_PAIRS, 1):
            if y < 1*inch:
                c.showPage()
                y = height - 1*inch
            
            c.setFont("Helvetica-Bold", 11)
            c.drawString(0.75*inch, y, f"Q{i}: {q}")
            y -= 0.3*inch
            c.setFont("Helvetica", 10)
            c.drawString(0.75*inch, y, f"A: {a}")
            y -= 0.5*inch
        
        c.save()
        print(f"Created: {filename}")
        return True
    except ImportError:
        print("Skipped PDF (install: pip install reportlab)")
        return False

def create_docx():
    """Create a DOCX file using python-docx."""
    try:
        from docx import Document
        from docx.shared import Pt
        
        filename = "current_affairs_qa.docx"
        doc = Document()
        
        # Title
        title = doc.add_heading("Current Affairs Questions and Answers - 2026", 0)
        
        # Q&A pairs
        for i, (q, a) in enumerate(QA_PAIRS, 1):
            doc.add_heading(f"Q{i}: {q}", level=2)
            p = doc.add_paragraph(f"Answer: {a}")
            doc.add_paragraph()  # Space
        
        doc.save(filename)
        print(f"Created: {filename}")
        return True
    except ImportError:
        print("Skipped DOCX (install: pip install python-docx)")
        return False

def create_pptx():
    """Create a PPTX file using python-pptx."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        filename = "current_affairs_qa.pptx"
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Current Affairs Q&A 2026"
        subtitle.text = "General Knowledge Quiz"
        
        # Q&A slides
        bullet_layout = prs.slide_layouts[1]
        for i, (q, a) in enumerate(QA_PAIRS[:5], 1):  # First 5 questions
            slide = prs.slides.add_slide(bullet_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            
            title_shape.text = f"Question {i}"
            tf = body_shape.text_frame
            tf.text = q
            p = tf.add_paragraph()
            p.text = f"Answer: {a}"
        
        prs.save(filename)
        print(f"Created: {filename}")
        return True
    except ImportError:
        print("Skipped PPTX (install: pip install python-pptx)")
        return False

def create_xlsx():
    """Create an XLSX file using openpyxl."""
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment
        
        filename = "current_affairs_qa.xlsx"
        wb = Workbook()
        ws = wb.active
        ws.title = "Q&A"
        
        # Headers
        ws['A1'] = 'Question Number'
        ws['B1'] = 'Question'
        ws['C1'] = 'Answer'
        
        # Style headers
        for cell in ws[1]:
            cell.font = Font(bold=True)
        
        # Add Q&A
        for i, (q, a) in enumerate(QA_PAIRS, 1):
            ws.cell(row=i+1, column=1, value=i)
            ws.cell(row=i+1, column=2, value=q)
            ws.cell(row=i+1, column=3, value=a)
        
        # Adjust column widths
        ws.column_dimensions['A'].width = 15
        ws.column_dimensions['B'].width = 45
        ws.column_dimensions['C'].width = 45
        
        wb.save(filename)
        print(f"Created: {filename}")
        return True
    except ImportError:
        print("Skipped XLSX (install: pip install openpyxl)")
        return False

def create_png():
    """Create a PNG image with text."""
    try:
        from PIL import Image, ImageDraw, ImageFont
        
        filename = "current_affairs_qa.png"
        
        # Create image
        img = Image.new('RGB', (800, 600), color='white')
        draw = ImageDraw.Draw(img)
        
        # Try to use a font, fallback to default
        try:
            font_title = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 24)
            font_text = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 14)
        except:
            font_title = ImageFont.load_default()
            font_text = ImageFont.load_default()
        
        # Title
        draw.text((50, 30), "Current Affairs Q&A 2026", fill='black', font=font_title)
        
        # Add some Q&A
        y = 80
        for i, (q, a) in enumerate(QA_PAIRS[:5], 1):
            draw.text((50, y), f"Q{i}: {q}", fill='navy', font=font_text)
            y += 25
            draw.text((50, y), f"A: {a}", fill='darkgreen', font=font_text)
            y += 35
        
        img.save(filename)
        print(f"Created: {filename}")
        return True
    except ImportError:
        print("Skipped PNG (install: pip install pillow)")
        return False

def create_jpeg():
    """Create a JPEG image with text."""
    try:
        from PIL import Image, ImageDraw, ImageFont
        
        filename = "current_affairs_qa.jpg"
        
        # Create image
        img = Image.new('RGB', (800, 600), color='#f0f0f0')
        draw = ImageDraw.Draw(img)
        
        try:
            font_title = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 24)
            font_text = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 14)
        except:
            font_title = ImageFont.load_default()
            font_text = ImageFont.load_default()
        
        draw.text((50, 30), "Current Affairs Quiz - 2026", fill='#333', font=font_title)
        
        y = 80
        for i, (q, a) in enumerate(QA_PAIRS[5:], 6):  # Questions 6-10
            draw.text((50, y), f"Q{i}: {q}", fill='#1a5276', font=font_text)
            y += 25
            draw.text((50, y), f"A: {a}", fill='#145a32', font=font_text)
            y += 35
        
        img.save(filename, 'JPEG', quality=85)
        print(f"Created: {filename}")
        return True
    except ImportError:
        print("Skipped JPEG (install: pip install pillow)")
        return False

def main():
    """Generate all binary test files."""
    print("Generating binary test files...\n")
    
    # Change to script directory
    script_dir = os.path.dirname(os.path.abspath(__file__))
    os.chdir(script_dir)
    
    created = []
    skipped = []
    
    if create_pdf(): created.append("PDF")
    else: skipped.append("PDF")
    
    if create_docx(): created.append("DOCX")
    else: skipped.append("DOCX")
    
    if create_pptx(): created.append("PPTX")
    else: skipped.append("PPTX")
    
    if create_xlsx(): created.append("XLSX")
    else: skipped.append("XLSX")
    
    if create_png(): created.append("PNG")
    else: skipped.append("PNG")
    
    if create_jpeg(): created.append("JPEG")
    else: skipped.append("JPEG")
    
    print(f"\n✅ Created: {', '.join(created) if created else 'None'}")
    if skipped:
        print(f"⚠️  Skipped (missing libraries): {', '.join(skipped)}")
        print("\nTo create all files, install required libraries:")
        print("  pip install reportlab python-docx python-pptx openpyxl pillow")

if __name__ == "__main__":
    main()

